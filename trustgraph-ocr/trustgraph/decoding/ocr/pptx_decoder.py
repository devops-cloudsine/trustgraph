"""
PPTX decoder: Extracts text, tables, and images from PPTX presentations using python-pptx.
Images are sent to vLLM for description. All content is output as structured text segments
suitable for RAG processing.

Key features:
- Uses python-pptx (free, open-source, no license required)
- Fallback to zipfile extraction when python-pptx fails (handles corrupt files)
- Concurrent vLLM requests for faster image processing
- Image deduplication by MD5 hash to avoid redundant LLM calls
- Robust error handling that continues processing even if parts fail
- Extracts text, tables, speaker notes, and images
- Images stored in MinIO for persistence

Reference: https://python-pptx.readthedocs.io/en/latest/api/presentation.html
"""

import asyncio
import base64
import concurrent.futures
import hashlib
import logging
import os
import re
import time
import urllib.parse
import zipfile
from io import BytesIO
from pathlib import Path
import requests
import json
from typing import Dict, List, Any, Optional, Set, Tuple

from ... schema import Document as TGDocument, TextDocument
from ... schema import Triples, Triple, Value, Metadata
from ... schema import EntityContext, EntityContexts
from ... base import FlowProcessor, ConsumerSpec, ProducerSpec
from .minio_storage import MinioStorage, get_minio_storage

# ---> rdf.py constants for image context triples
IMAGE_CONTEXT = "http://trustgraph.ai/ns/image-context"
IMAGE_SOURCE = "http://trustgraph.ai/ns/image-source"
TRUSTGRAPH_ENTITIES = "http://trustgraph.ai/e/"

# Module logger
logger = logging.getLogger(__name__)

# Try to import python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.shapes.picture import Picture
    from pptx.shapes.graphfrm import GraphicFrame
    from pptx.table import Table
    PPTX_AVAILABLE = True
except ImportError as e:
    PPTX_AVAILABLE = False
    Presentation = None
    logger.warning(f"python-pptx not available: {e}")

# Try to import PIL for image handling
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

# Enable DEBUG level for this module when VLLM_LOGGING_LEVEL=DEBUG is set
if os.environ.get("VLLM_LOGGING_LEVEL", "").upper() == "DEBUG":
    logger.setLevel(logging.DEBUG)

default_ident = "pptx-decoder"

# Performance tuning constants
MAX_CONCURRENT_VLLM_REQUESTS = 4  # Max parallel vLLM requests
VLLM_TIMEOUT_SECONDS = 180  # Timeout for vLLM requests
MAX_IMAGE_SIZE_MB = 10  # Skip images larger than this
VLLM_MAX_RETRIES = 3  # Max retries for vLLM requests
VLLM_RETRY_DELAY_SECONDS = 2  # Delay between retries


class Processor(FlowProcessor):

    def __init__(self, **params):
        id = params.get("id", default_ident)
        self.vllm_api_url = params.get(
            "vllm_api_url",
            "http://vllm:8000/v1/chat/completions"
        )
        self.vllm_model = params.get(
            "vllm_model",
            "Qwen/Qwen3-VL-4B-Instruct"
        )
        
        # Initialize MinIO storage for image persistence
        self.minio_storage = get_minio_storage()

        super(Processor, self).__init__(
            **params | {"id": id}
        )

        self.register_specification(
            ConsumerSpec(
                name="input",
                schema=TGDocument,
                handler=self.on_message,
            )
        )

        self.register_specification(
            ProducerSpec(
                name="output",
                schema=TextDocument,
            )
        )

        # ---> on_message > [triples output] > emit IMAGE_CONTEXT/IMAGE_SOURCE triples for Graph RAG
        self.register_specification(
            ProducerSpec(
                name="triples",
                schema=Triples,
            )
        )

        # ---> on_message > [entity-contexts output] > emit EntityContexts for graph embeddings (enables graph-retrieval)
        self.register_specification(
            ProducerSpec(
                name="entity-contexts",
                schema=EntityContexts,
            )
        )

        logger.info("PPTX decoder initialized (python-pptx + zipfile fallback, concurrent vLLM)")

    # // ---> Pulsar consumer(input) > [on_message] > extract slides/text/images, vLLM describe -> flow('output')
    async def on_message(self, msg, consumer, flow):
        """
        Main message handler with top-level exception handling.
        Pulsar will redeliver the message if this method raises an exception.
        """
        v = None
        try:
            logger.info("PPTX message received for extraction")

            v = msg.value()
            doc_id = v.metadata.id
            logger.info(f"Processing {doc_id}...")

            blob = base64.b64decode(v.data)

            # Check if this is a PPTX file
            content_type = getattr(v, "content_type", None)
            if not self._is_pptx(blob, content_type):
                logger.info(f"Skipping non-PPTX file: {doc_id} (content_type: {content_type})")
                return

            # Initialize MinIO storage
            if not self.minio_storage.initialize():
                logger.error("Failed to initialize MinIO storage - cannot process PPTX")
                return

            # Save original PPTX document to MinIO
            pptx_filename = f"{self._safe_id(doc_id)}.pptx"
            doc_object_name = self.minio_storage.save_document(
                doc_id=doc_id,
                filename=pptx_filename,
                document_data=blob,
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            if doc_object_name:
                logger.info(f"Saved original PPTX to MinIO: {doc_object_name}")
            else:
                logger.warning("Failed to save original PPTX to MinIO, continuing with processing")

            # Track processed image hashes for deduplication
            processed_hashes: Set[str] = set()
            all_images: List[Dict[str, Any]] = []
            structured_data = None
            pptx_success = False

            # Try python-pptx extraction
            if PPTX_AVAILABLE:
                pptx_success, structured_data, all_images = self._try_pptx_extraction(
                    blob, doc_id, processed_hashes
                )

            # Fallback to zipfile extraction if python-pptx failed
            if not pptx_success:
                logger.info("Using zipfile fallback extraction")
                fallback_images = self._extract_images_from_zip(blob, doc_id, processed_hashes)
                for ix, img_info in enumerate(fallback_images):
                    all_images.append({
                        "data": img_info["data"],
                        "content_type": img_info["content_type"],
                        "context": f"Image {ix + 1} ({img_info['original_name']})",
                        "slide_number": None,
                        "shape_index": None
                    })
                structured_data = {
                    "presentation": {
                        "total_slides": 0,
                        "slides": [],
                        "note": "Extracted via zipfile fallback"
                    }
                }
                logger.info(f"Fallback: Found {len(fallback_images)} images")

            # Send presentation header first
            if structured_data:
                # Get filename from metadata or path
                filename = getattr(v.metadata, 'id', 'presentation') + '.pptx'
                header = self._format_presentation_header(structured_data, filename)
                if header.strip():
                    logger.info(f"Sending presentation header to RAG ({len(header)} chars)")
                    logger.debug(f"Header content:\n{header[:500]}{'...' if len(header) > 500 else ''}")
                    r = TextDocument(
                        metadata=v.metadata,
                        text=header.encode("utf-8"),
                    )
                    await flow("output").send(r)

            # Send text content from each slide
            if structured_data and structured_data["presentation"]["slides"]:
                for slide_data in structured_data["presentation"]["slides"]:
                    slide_text = self._format_slide_text(slide_data)
                    if slide_text.strip():
                        slide_num = slide_data.get("slide_number", "?")
                        logger.info(f"Sending slide {slide_num} content to RAG ({len(slide_text)} chars)")
                        logger.debug(f"Slide {slide_num} content:\n{slide_text[:500]}{'...' if len(slide_text) > 500 else ''}")
                        r = TextDocument(
                            metadata=v.metadata,
                            text=slide_text.encode("utf-8"),
                        )
                        await flow("output").send(r)

            # Process images with concurrent vLLM requests
            # Collect triples and entity contexts for Graph RAG
            image_triples = []
            entity_contexts = []  # For graph embeddings
            doc_uri = TRUSTGRAPH_ENTITIES + urllib.parse.quote(doc_id)
            
            if all_images:
                logger.info(f"Processing {len(all_images)} unique images with concurrent vLLM requests")
                image_descriptions = await self._describe_images_concurrent_bytes(all_images, doc_id)

                # Send image descriptions and collect triples + entity contexts
                for idx, (img_info, description) in enumerate(zip(all_images, image_descriptions)):
                    if description and description.strip():
                        image_text = f"{img_info['context']}\nDescription:\n{description}\n"
                        logger.info(f"Sending image {idx+1}/{len(all_images)} description to RAG ({len(image_text)} chars)")
                        logger.debug(f"Image description:\n{image_text[:300]}{'...' if len(image_text) > 300 else ''}")
                        r = TextDocument(
                            metadata=v.metadata,
                            text=image_text.encode("utf-8"),
                        )
                        await flow("output").send(r)
                        
                        # ---> Create IMAGE_CONTEXT and IMAGE_SOURCE triples for Graph RAG
                        slide_num = img_info.get('slide_number', idx + 1)
                        image_uri = f"{doc_uri}/slide_{slide_num}/image_{idx + 1}"
                        image_uri_value = Value(value=image_uri, is_uri=True)
                        
                        # Add IMAGE_CONTEXT triple
                        image_triples.append(Triple(
                            s=image_uri_value,
                            p=Value(value=IMAGE_CONTEXT, is_uri=True),
                            o=Value(value=description, is_uri=False),
                        ))
                        
                        # Add IMAGE_SOURCE triple
                        image_path = f"minio://ocr-images/{self._safe_id(doc_id)}/image_{idx}.png"
                        image_triples.append(Triple(
                            s=image_uri_value,
                            p=Value(value=IMAGE_SOURCE, is_uri=True),
                            o=Value(value=image_path, is_uri=False),
                        ))
                        
                        # ---> Create EntityContext for graph embeddings (enables graph-retrieval to find this entity)
                        entity_contexts.append(EntityContext(
                            entity=image_uri_value,
                            context=description,
                        ))

            # ---> on_message > [triples flow] > emit IMAGE_CONTEXT/IMAGE_SOURCE triples for Graph RAG
            if image_triples:
                triples_msg = Triples(
                    metadata=Metadata(
                        id=doc_id,
                        metadata=[],
                        user=v.metadata.user if hasattr(v.metadata, 'user') and v.metadata.user else "trustgraph",
                        collection=v.metadata.collection if hasattr(v.metadata, 'collection') and v.metadata.collection else "default",
                    ),
                    triples=image_triples,
                )
                await flow("triples").send(triples_msg)
                logger.info(f"Emitted {len(image_triples)} image context triples for document {doc_id}")

            # ---> on_message > [entity-contexts flow] > emit EntityContexts for graph embeddings (enables graph-retrieval)
            if entity_contexts:
                entity_contexts_msg = EntityContexts(
                    metadata=Metadata(
                        id=doc_id,
                        metadata=[],
                        user=v.metadata.user if hasattr(v.metadata, 'user') and v.metadata.user else "trustgraph",
                        collection=v.metadata.collection if hasattr(v.metadata, 'collection') and v.metadata.collection else "default",
                    ),
                    entities=entity_contexts,
                )
                await flow("entity-contexts").send(entity_contexts_msg)
                logger.info(f"Emitted {len(entity_contexts)} entity contexts for graph embeddings for document {doc_id}")

            logger.info(f"PPTX extraction complete: {len(all_images)} images processed")

        except Exception as e:
            doc_id = v.metadata.id if v and hasattr(v, 'metadata') else "unknown"
            logger.error(f"Fatal error processing PPTX {doc_id}: {type(e).__name__}: {e}")
            # Re-raise to trigger Pulsar redelivery for transient errors
            raise

    # // ---> on_message > [_try_pptx_extraction] > extract using python-pptx with fallback
    def _try_pptx_extraction(
        self, blob: bytes, doc_id: str, processed_hashes: Set[str]
    ) -> Tuple[bool, Optional[Dict], List[Dict]]:
        """
        Try to extract content using python-pptx.
        Falls back to zipfile if python-pptx fails completely.
        Images are stored in memory with bytes data for MinIO upload.
        """
        all_images = []
        structured_data = None

        try:
            # Load from bytes (in-memory)
            logger.debug("Trying python-pptx load from bytes")
            prs = Presentation(BytesIO(blob))
            logger.info(f"python-pptx loaded successfully: {len(prs.slides)} slides")

            # Extract all content
            structured_data = self._extract_structured_pptx(prs, doc_id, processed_hashes)
            logger.info(f"Extracted {structured_data['presentation']['total_slides']} slides")

            # Collect slide images with bytes data
            for slide_data in structured_data["presentation"]["slides"]:
                for img_info in slide_data.get("images", []):
                    all_images.append({
                        "data": img_info["data"],
                        "content_type": img_info["content_type"],
                        "context": f"Slide {slide_data['slide_number']} - {img_info['filename']}",
                        "slide_number": slide_data["slide_number"],
                        "shape_index": img_info.get("shape_index")
                    })

            return True, structured_data, all_images

        except Exception as e:
            logger.warning(f"python-pptx load failed: {e}")
            return False, None, []

    # // ---> on_message > [_process_with_fallback] > handle fallback when extraction fails
    async def _process_with_fallback(self, blob: bytes, doc_id: str, v, flow):
        """Fallback processing when python-pptx fails."""
        processed_hashes: Set[str] = set()
        images = self._extract_images_from_zip(blob, doc_id, processed_hashes)

        if images:
            logger.info(f"Fallback: Extracted {len(images)} images")
            all_images = [
                {
                    "data": img["data"],
                    "content_type": img["content_type"],
                    "context": f"Image {i+1} ({img['original_name']})",
                    "slide_number": None,
                    "shape_index": None
                }
                for i, img in enumerate(images)
            ]
            descriptions = await self._describe_images_concurrent_bytes(all_images, doc_id)

            for img_info, desc in zip(all_images, descriptions):
                if desc:
                    r = TextDocument(
                        metadata=v.metadata,
                        text=f"{img_info['context']}\nDescription:\n{desc}\n".encode("utf-8")
                    )
                    await flow("output").send(r)
        else:
            logger.warning("Fallback extraction found no images")

    @staticmethod
    def add_args(parser):
        parser.add_argument(
            '--vllm-api-url',
            default='http://vllm-vision-server:8000/v1/chat/completions',
            help='vLLM API URL for image descriptions (default: http://vllm:8000/v1/chat/completions)'
        )
        parser.add_argument(
            '--vllm-model',
            default='google/gemma-3-4b-it',
            help='vLLM model name (default: Qwen/Qwen3-VL-4B-Instruct)'
        )
        FlowProcessor.add_args(parser)

    # // ---> [_safe_id] > sanitize directory name
    def _safe_id(self, value: str) -> str:
        if value is None:
            return "unknown"
        return re.sub(r"[^A-Za-z0-9._-]+", "_", str(value))

    # // ---> [_is_pptx] > check if blob is a PPTX file
    def _is_pptx(self, blob: bytes, content_type: Optional[str] = None) -> bool:
        """
        Check if blob is a PPTX file by verifying actual file content (magic bytes + ZIP structure).
        ALWAYS checks magic bytes first, regardless of content_type.
        """
        # PPTX files are ZIP archives - ALWAYS check magic bytes first
        if len(blob) < 4 or blob[:4] != b'PK\x03\x04':
            if content_type:
                logger.debug(f"File has content_type={content_type} but is not a ZIP archive (not PPTX)")
            return False

        # Inspect ZIP contents to distinguish PPTX from DOCX/XLSX
        # PPTX has ppt/ directory, DOCX has word/, XLSX has xl/
        try:
            with zipfile.ZipFile(BytesIO(blob), 'r') as zf:
                names = zf.namelist()
                # Check for PPTX-specific paths
                has_ppt = any(n.startswith('ppt/') for n in names)
                has_content_types = '[Content_Types].xml' in names
                # Exclude DOCX and XLSX
                has_word = any(n.startswith('word/') for n in names)
                has_xl = any(n.startswith('xl/') for n in names)

                if has_ppt and has_content_types and not has_word and not has_xl:
                    return True
        except Exception:
            pass

        return False

    # // ---> [_compute_image_hash] > MD5 hash for deduplication
    def _compute_image_hash(self, image_data: bytes) -> str:
        return hashlib.md5(image_data).hexdigest()

    # // ---> [_extract_images_from_zip] > fallback extraction via zipfile
    def _extract_images_from_zip(
        self, blob: bytes, doc_id: str, processed_hashes: Set[str]
    ) -> List[Dict[str, Any]]:
        """
        Extract images directly from PPTX ZIP archive.
        Returns list of dicts with image data, content_type, and original_name.
        Works even when python-pptx can't parse the file.
        """
        extracted = []
        image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'}
        # EMF/WMF are Windows metafiles - we'll try to handle them
        metafile_extensions = {'.emf', '.wmf'}
        
        ext_to_content_type = {
            '.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
            '.gif': 'image/gif', '.bmp': 'image/bmp', '.tiff': 'image/tiff',
            '.webp': 'image/webp'
        }

        try:
            with zipfile.ZipFile(BytesIO(blob), 'r') as zf:
                for name in zf.namelist():
                    # PPTX stores media in ppt/media/
                    if 'media/' not in name.lower():
                        continue

                    ext = Path(name).suffix.lower()
                    original_name = Path(name).name

                    if ext in image_extensions or ext in metafile_extensions:
                        try:
                            image_data = zf.read(name)

                            # Skip very large images
                            if len(image_data) > MAX_IMAGE_SIZE_MB * 1024 * 1024:
                                logger.debug(f"Skipping large image: {name}")
                                continue

                            img_hash = self._compute_image_hash(image_data)
                            if img_hash in processed_hashes:
                                continue
                            processed_hashes.add(img_hash)

                            # Try to convert EMF/WMF if PIL is available
                            if ext in metafile_extensions and PIL_AVAILABLE:
                                converted = self._try_convert_metafile(image_data)
                                if converted:
                                    image_data = converted
                                    ext = '.png'
                                else:
                                    # Skip unconvertible metafiles
                                    logger.debug(f"Skipping unconvertible metafile: {name}")
                                    continue

                            content_type = ext_to_content_type.get(ext, 'image/png')
                            extracted.append({
                                "data": image_data,
                                "content_type": content_type,
                                "original_name": original_name,
                                "filename": f"zip_{len(extracted)}{ext}"
                            })
                            logger.debug(f"Extracted from ZIP: {original_name}")

                        except Exception as e:
                            logger.debug(f"Failed to extract {name}: {e}")

        except zipfile.BadZipFile:
            logger.warning("Not a valid ZIP archive")
        except Exception as e:
            logger.warning(f"ZIP extraction failed: {e}")

        return extracted

    # // ---> [_try_convert_metafile] > convert EMF/WMF to PNG
    def _try_convert_metafile(self, data: bytes) -> Optional[bytes]:
        """Try to convert EMF/WMF metafile to PNG using PIL."""
        if not PIL_AVAILABLE:
            return None
        try:
            img = Image.open(BytesIO(data))
            output = BytesIO()
            img.save(output, format='PNG')
            return output.getvalue()
        except Exception:
            return None

    # // ---> [_extract_structured_pptx] > parse PPTX with python-pptx
    def _extract_structured_pptx(
        self, prs, doc_id: str, processed_hashes: Set[str]
    ) -> Dict[str, Any]:
        """Extract all content from PPTX using python-pptx. Images stored as bytes."""
        structured_data = {
            "presentation": {
                "total_slides": len(prs.slides),
                "slides": [],
                "properties": {}
            }
        }

        # Extract core properties
        try:
            props = prs.core_properties
            structured_data["presentation"]["properties"] = {
                "title": props.title or "",
                "author": props.author or "",
                "subject": props.subject or "",
                "keywords": props.keywords or "",
                "comments": props.comments or "",
            }
        except Exception as e:
            logger.debug(f"Failed to extract core properties: {e}")

        image_counter = 0

        # Process each slide
        for slide_idx, slide in enumerate(prs.slides):
            try:
                slide_data = self._extract_slide_data(
                    slide, slide_idx, processed_hashes, image_counter
                )
                image_counter += len(slide_data.get("images", []))
                structured_data["presentation"]["slides"].append(slide_data)
            except Exception as e:
                logger.warning(f"Failed to process slide {slide_idx + 1}: {e}")
                structured_data["presentation"]["slides"].append({
                    "slide_number": slide_idx + 1,
                    "error": str(e),
                    "shapes": [], "text_content": [], "images": [], "tables": []
                })

        return structured_data

    # // ---> [_extract_slide_data] > extract all content from a single slide
    def _extract_slide_data(
        self, slide, slide_idx: int,
        processed_hashes: Set[str], image_counter: int
    ) -> Dict[str, Any]:
        """Extract all content from a single slide. Images stored as bytes."""
        slide_data = {
            "slide_number": slide_idx + 1,
            "slide_id": None,
            "layout": None,
            "shapes": [],
            "text_content": [],
            "images": [],
            "tables": [],
            "notes": None
        }

        # Get slide layout name
        try:
            if slide.slide_layout:
                slide_data["layout"] = slide.slide_layout.name
        except Exception:
            pass

        # Extract speaker notes
        try:
            if slide.has_notes_slide and slide.notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text:
                    slide_data["notes"] = notes_frame.text.strip()
        except Exception as e:
            logger.debug(f"Notes extraction failed: {e}")

        # Process all shapes
        local_img_counter = 0
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                shape_data, new_images = self._extract_shape_data(
                    shape, shape_idx, slide_idx,
                    processed_hashes, image_counter + local_img_counter
                )
                slide_data["shapes"].append(shape_data)

                # Add text content
                if shape_data.get("text"):
                    slide_data["text_content"].append({
                        "shape_index": shape_idx,
                        "shape_name": shape_data.get("name"),
                        "text": shape_data["text"]
                    })

                # Add images
                for img in new_images:
                    slide_data["images"].append(img)
                    local_img_counter += 1

                # Add table
                if shape_data.get("table_data"):
                    slide_data["tables"].append({
                        "shape_index": shape_idx,
                        "rows": len(shape_data["table_data"]),
                        "columns": len(shape_data["table_data"][0]) if shape_data["table_data"] else 0,
                        "data": shape_data["table_data"]
                    })

            except Exception as e:
                logger.debug(f"Shape {shape_idx} extraction failed: {e}")

        return slide_data

    # // ---> [_extract_shape_data] > extract content from a single shape
    def _extract_shape_data(
        self, shape, shape_idx: int, slide_idx: int,
        processed_hashes: Set[str], image_counter: int
    ) -> Tuple[Dict[str, Any], List[Dict]]:
        """Extract all content from a shape. Returns (shape_data, list of images with bytes)."""
        shape_data = {
            "shape_index": shape_idx,
            "shape_type": None,
            "name": None,
            "text": None,
            "table_data": None,
            "image_ref": None  # Will be set if shape contains an image
        }
        images = []

        # Get shape name and type
        try:
            shape_data["name"] = shape.name
        except Exception:
            pass

        try:
            shape_data["shape_type"] = str(shape.shape_type)
        except Exception:
            pass

        # Extract text from text frame
        try:
            if shape.has_text_frame:
                text_parts = []
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        text_parts.append(para_text)
                if text_parts:
                    shape_data["text"] = "\n".join(text_parts)
        except Exception as e:
            logger.debug(f"Text extraction failed for shape {shape_idx}: {e}")

        # Extract image if shape is a picture
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_info = self._extract_picture(
                    shape, slide_idx, processed_hashes, image_counter
                )
                if img_info:
                    images.append({
                        "shape_index": shape_idx,
                        "filename": img_info["filename"],
                        "data": img_info["data"],
                        "content_type": img_info["content_type"]
                    })
                    # Add image reference to shape data for shape summary
                    shape_data["image_ref"] = img_info["filename"]
        except Exception as e:
            logger.debug(f"Picture extraction failed for shape {shape_idx}: {e}")

        # Extract table data if shape is a table
        try:
            if shape.has_table:
                table_data = self._extract_table_data(shape.table)
                if table_data:
                    shape_data["table_data"] = table_data
        except Exception as e:
            logger.debug(f"Table extraction failed for shape {shape_idx}: {e}")

        # Check for grouped shapes
        try:
            if hasattr(shape, 'shapes'):
                # This is a group shape, extract from children
                for child_idx, child_shape in enumerate(shape.shapes):
                    try:
                        child_data, child_images = self._extract_shape_data(
                            child_shape, f"{shape_idx}_{child_idx}", slide_idx,
                            processed_hashes, image_counter + len(images)
                        )
                        # Merge child text
                        if child_data.get("text"):
                            if shape_data["text"]:
                                shape_data["text"] += "\n" + child_data["text"]
                            else:
                                shape_data["text"] = child_data["text"]
                        images.extend(child_images)
                    except Exception:
                        pass
        except Exception:
            pass

        return shape_data, images

    # // ---> [_extract_picture] > extract image from Picture shape
    def _extract_picture(
        self, shape, slide_idx: int,
        processed_hashes: Set[str], counter: int
    ) -> Optional[Dict[str, Any]]:
        """Extract image from a Picture shape with deduplication. Returns bytes data."""
        try:
            # Get image blob from shape
            image = shape.image
            image_bytes = image.blob

            # Check size
            if len(image_bytes) > MAX_IMAGE_SIZE_MB * 1024 * 1024:
                logger.debug(f"Skipping large image in shape")
                return None

            # Compute hash for deduplication
            img_hash = self._compute_image_hash(image_bytes)
            if img_hash in processed_hashes:
                logger.debug(f"Skipping duplicate image")
                return None
            processed_hashes.add(img_hash)

            # Determine content type and extension
            content_type = image.content_type
            ext_map = {
                'image/png': '.png',
                'image/jpeg': '.jpg',
                'image/gif': '.gif',
                'image/bmp': '.bmp',
                'image/tiff': '.tiff',
                'image/webp': '.webp',
                'image/x-emf': '.emf',
                'image/x-wmf': '.wmf',
            }
            ext = ext_map.get(content_type, '.png')

            # Handle metafiles
            if ext in ['.emf', '.wmf'] and PIL_AVAILABLE:
                converted = self._try_convert_metafile(image_bytes)
                if converted:
                    image_bytes = converted
                    content_type = 'image/png'
                    ext = '.png'
                else:
                    # Skip unconvertible
                    return None

            filename = f"slide_{slide_idx + 1}_img_{counter}{ext}"

            logger.debug(f"Extracted picture: {filename}")
            return {
                "filename": filename,
                "data": image_bytes,
                "content_type": content_type
            }

        except Exception as e:
            logger.debug(f"Failed to extract picture: {e}")
            return None

    # // ---> [_extract_table_data] > extract table cell data
    def _extract_table_data(self, table) -> Optional[List[List[str]]]:
        """Extract table data as 2D list of strings."""
        try:
            rows_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    try:
                        cell_text = cell.text.strip() if cell.text else ""
                        row_data.append(cell_text)
                    except Exception:
                        row_data.append("")
                rows_data.append(row_data)
            return rows_data if rows_data else None
        except Exception as e:
            logger.debug(f"Table extraction failed: {e}")
            return None

    # // ---> [_format_presentation_header] > format presentation header
    def _format_presentation_header(self, structured_data: Dict, filename: str) -> str:
        """Format presentation header with metadata."""
        parts = []
        pres = structured_data.get("presentation", {})

        parts.append(f"# PowerPoint Presentation: {filename}")
        parts.append("")
        parts.append(f"Total Slides: {pres.get('total_slides', 0)}")

        # Include core properties if available
        props = pres.get("properties", {})
        if props.get("title"):
            parts.append(f"Title: {props['title']}")
        if props.get("author"):
            parts.append(f"Author: {props['author']}")
        if props.get("subject"):
            parts.append(f"Subject: {props['subject']}")

        parts.append("")
        return "\n".join(parts)

    # // ---> [_format_slide_text] > format slide data as text
    def _format_slide_text(self, slide_data: Dict) -> str:
        """Format slide data as structured text for RAG (matching expected format)."""
        parts = []
        slide_num = slide_data["slide_number"]

        # Slide header
        parts.append(f"{'='*80}")
        parts.append("")
        parts.append(f"SLIDE {slide_num}")
        if slide_data.get("layout"):
            parts.append(f"Layout: {slide_data['layout']}")
        parts.append("")
        parts.append(f"{'='*80}")
        parts.append("")

        # Text content section
        if slide_data.get("text_content"):
            parts.append("## Text Content:")
            parts.append("")
            for item in slide_data["text_content"]:
                text = item.get('text', '')
                if text:
                    parts.append(text)
                    parts.append("")

        # Tables section
        if slide_data.get("tables"):
            parts.append("## Tables:")
            parts.append("")
            for idx, table in enumerate(slide_data["tables"]):
                parts.append(f"Table {idx + 1} ({table['rows']} rows × {table['columns']} columns):")
                for row in table.get("data", []):
                    parts.append(" | ".join(str(c) if c else "" for c in row))
                parts.append("")

        # Images section
        if slide_data.get("images"):
            parts.append("## Images:")
            parts.append("")
            for img in slide_data["images"]:
                shape_idx = img.get("shape_index", "N/A")
                filename = img.get("filename", "unknown")
                parts.append(f"- Image: {filename} (Shape Index: {shape_idx})")
            parts.append("")

        # Speaker Notes section
        if slide_data.get("notes"):
            parts.append("## Speaker Notes:")
            parts.append("")
            parts.append(slide_data["notes"])
            parts.append("")

        # Shape Summary section
        shapes = slide_data.get("shapes", [])
        if shapes:
            parts.append(f"## Shape Summary: {len(shapes)} shapes")
            parts.append("")
            for shape in shapes:
                shape_idx = shape.get("shape_index", "?")
                shape_type = shape.get("shape_type", "Unknown")
                shape_name = shape.get("name", "")

                # Build shape info line
                shape_info = f"  - Shape {shape_idx}: {shape_type}"
                if shape_name:
                    shape_info += f" (Name: {shape_name})"

                # Add text preview if available
                if shape.get("text"):
                    text_preview = shape["text"][:50]
                    if len(shape["text"]) > 50:
                        text_preview += "..."
                    shape_info += f" - Text: {text_preview}"

                # Add image reference if available
                if shape.get("image_ref"):
                    shape_info += f" - Image: {shape['image_ref']}"

                parts.append(shape_info)
            parts.append("")

        return "\n".join(parts)

    # // ---> [_bytes_to_base64_data_url] > convert image bytes to base64 data URL
    def _bytes_to_base64_data_url(self, image_bytes: bytes, content_type: str = "image/png") -> str:
        """Convert image bytes to base64 data URL for vLLM."""
        b64 = base64.b64encode(image_bytes).decode("utf-8")
        return f"data:{content_type};base64,{b64}"

    # // ---> [_describe_images_concurrent_bytes] > concurrent vLLM requests with bytes
    async def _describe_images_concurrent_bytes(self, images: List[Dict], doc_id: str) -> List[str]:
        """
        Describe images using concurrent vLLM requests.
        Also saves images to MinIO storage.
        Processes up to MAX_CONCURRENT_VLLM_REQUESTS in parallel.
        """
        if not images:
            return []

        # Filter valid images (must have data)
        valid_images = [img for img in images if img.get("data")]
        if not valid_images:
            return [""] * len(images)

        logger.info(f"Processing {len(valid_images)} images concurrently")

        # Use get_running_loop() for Python 3.10+ compatibility
        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            loop = asyncio.get_event_loop()

        with concurrent.futures.ThreadPoolExecutor(max_workers=MAX_CONCURRENT_VLLM_REQUESTS) as executor:
            futures = [
                loop.run_in_executor(
                    executor,
                    self._describe_and_save_image,
                    img["data"],
                    img.get("content_type", "image/png"),
                    doc_id,
                    idx
                )
                for idx, img in enumerate(valid_images)
            ]
            results = await asyncio.gather(*futures, return_exceptions=True)

        # Process results - maintain order matching input images
        descriptions = []
        for result in results:
            if isinstance(result, Exception):
                logger.error(f"Image description failed: {result}")
                descriptions.append("")
            else:
                descriptions.append(result or "")

        return descriptions

    # // ---> [_describe_and_save_image] > save to MinIO and describe via vLLM with retry
    def _describe_and_save_image(
        self, image_bytes: bytes, content_type: str, doc_id: str, image_idx: int
    ) -> str:
        """
        Save image to MinIO and describe using vLLM with retry logic.
        Retries on connection errors, timeouts, and 5xx server errors.
        """
        # Determine extension from content type
        ext_map = {
            'image/png': '.png', 'image/jpeg': '.jpg', 'image/gif': '.gif',
            'image/bmp': '.bmp', 'image/webp': '.webp'
        }
        ext = ext_map.get(content_type, '.png')
        filename = f"image_{image_idx}{ext}"
        
        # Save to MinIO
        object_name = self.minio_storage.save_image(
            doc_id=doc_id,
            filename=filename,
            image_data=image_bytes,
            content_type=content_type
        )
        if object_name:
            logger.debug(f"Saved image to MinIO: {object_name}")
        
        # Describe via vLLM
        image_url = self._bytes_to_base64_data_url(image_bytes, content_type)

        payload = {
            "model": self.vllm_model,
            "messages": [{
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            "Describe this image from a presentation slide in detail. "
                            "Focus on: 1) Any text or labels visible, 2) Charts/graphs and their data, "
                            "3) Diagrams and their relationships, 4) Key visual elements. "
                            "Be concise but comprehensive."
                        )
                    },
                    {"type": "image_url", "image_url": {"url": image_url}}
                ]
            }]
        }

        last_error = None
        for attempt in range(VLLM_MAX_RETRIES):
            try:
                resp = requests.post(
                    self.vllm_api_url,
                    json=payload,
                    timeout=VLLM_TIMEOUT_SECONDS
                )

                # Success
                if resp.status_code == 200:
                    data = resp.json()
                    choice = (data.get("choices") or [{}])[0]
                    content = (choice.get("message") or {}).get("content", "")

                    if isinstance(content, str):
                        logger.info(f"vLLM image description: {content.strip()}")
                        return content.strip()
                    if isinstance(content, list):
                        return "\n".join(c.get("text", "") for c in content if c.get("type") == "text")
                    return ""

                # Client errors (4xx) - don't retry
                if 400 <= resp.status_code < 500:
                    logger.error(f"vLLM client error HTTP {resp.status_code}: {resp.text[:200]}")
                    return f"Description unavailable (HTTP {resp.status_code})"

                # Server errors (5xx) - retry
                last_error = f"HTTP {resp.status_code}"
                logger.warning(f"vLLM server error {resp.status_code}, attempt {attempt + 1}/{VLLM_MAX_RETRIES}")

            except requests.exceptions.Timeout:
                last_error = "timeout"
                logger.warning(f"vLLM timeout, attempt {attempt + 1}/{VLLM_MAX_RETRIES}")

            except requests.exceptions.ConnectionError as e:
                last_error = f"connection error: {e}"
                logger.warning(f"vLLM connection error, attempt {attempt + 1}/{VLLM_MAX_RETRIES}: {e}")

            except Exception as e:
                last_error = str(e)
                logger.error(f"vLLM unexpected error: {type(e).__name__}: {e}")
                # Don't retry on unexpected errors
                break

            # Wait before retry (except on last attempt)
            if attempt < VLLM_MAX_RETRIES - 1:
                time.sleep(VLLM_RETRY_DELAY_SECONDS * (attempt + 1))  # Exponential backoff

        logger.error(f"vLLM failed after {VLLM_MAX_RETRIES} attempts: {last_error}")
        return f"Description unavailable ({last_error})"


def run():
    Processor.launch(default_ident, __doc__)
